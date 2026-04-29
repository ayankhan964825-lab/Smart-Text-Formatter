from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── Colours from reference PPT ──
BG    = RGBColor(0xF4,0xF1,0xE9)   # cream
BLUE  = RGBColor(0x51,0x70,0xFF)
CORAL = RGBColor(0xFF,0x6F,0x5E)
BLACK = RGBColor(0x00,0x00,0x00)
WHITE = RGBColor(0xFF,0xFF,0xFF)
LTBLUE= RGBColor(0xE8,0xEC,0xFF)   # formula box fill
DARKBG= RGBColor(0x18,0x18,0x2E)

TF = "Cormorant Garamond Bold"
BF = "Canva Sans"
MF = "Cambria Math"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Fade transition XML ──
FADE_XML = (
    '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    ' spd="med" advClick="1"><p:fade/></p:transition>'
)
def add_transition(slide):
    slide._element.append(etree.fromstring(FADE_XML))

# ── Wipe-in animation for a shape ──
ANIM_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
def add_fade_anim(slide, shape, delay_ms=0):
    """Add a simple fade-in entrance animation to a shape."""
    spid = shape.shape_id
    timing_xml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" presetID="10" presetClass="entr" presetSubtype="0"
                               fill="hold" grpId="0" nodeType="clickEffect" dur="500">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:set><p:cBhvr>
                              <p:cTn id="5" dur="1" fill="hold"/>
                              <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                              <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                            </p:cBhvr>
                              <p:to><p:strVal val="visible"/></p:to>
                            </p:set>
                            <p:animEffect transition="in" filter="fade">
                              <p:cBhvr><p:cTn id="6" dur="500"/>
                                <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                              </p:cBhvr>
                            </p:animEffect>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst><p:cond evt="onPrevClick" delay="0"><p:tn/></p:cond></p:prevCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
  <p:bldLst/>
</p:timing>'''
    try:
        slide._element.append(etree.fromstring(timing_xml))
    except:
        pass  # skip if timing already exists

# ── Helpers ──
def new_slide(bg_color=None):
    sl = prs.slides.add_slide(BLANK)
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color or BG
    add_transition(sl)
    return sl

def tb(sl, x, y, w, h):
    return sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))

def rn(tf, txt, font=BF, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(4)
    r = p.add_run()
    r.text = txt
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    return p

def heading_bar(sl, title):
    """Blue gradient title + coral underline bar."""
    t = tb(sl, 0.4, 0.2, 12.5, 1.0)
    tf = t.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = TF; run.font.size = Pt(38); run.font.bold = True
    run.font.color.rgb = BLUE
    # coral accent bar
    bar = sl.shapes.add_shape(1, Inches(0.4), Inches(1.15), Inches(12.5), Pt(4))
    bar.fill.solid(); bar.fill.fore_color.rgb = CORAL
    bar.line.fill.background()
    return t

def bullet_box(sl, items, x=0.4, y=1.35, w=12.5, h=5.9, size=17, numbered=False):
    t = tb(sl, x, y, w, h)
    tf = t.text_frame; tf.word_wrap = True
    first = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT; p.space_before = Pt(6)
        if numbered:
            nr = p.add_run(); nr.text = f"{i+1}.  "
            nr.font.name = BF; nr.font.size = Pt(size)
            nr.font.bold = True; nr.font.color.rgb = CORAL
        r = p.add_run(); r.text = item
        r.font.name = BF; r.font.size = Pt(size); r.font.color.rgb = BLACK
    return t

def formula_box(sl, lines, x=1.5, y=1.4, w=10.0, h=2.0):
    t = tb(sl, x, y, w, h)
    t.fill.solid(); t.fill.fore_color.rgb = LTBLUE
    tf = t.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(8)
        r = p.add_run(); r.text = line
        r.font.name = MF; r.font.size = Pt(20)
        r.font.bold = True; r.font.color.rgb = BLUE
    return t

# ══════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════
sl = new_slide()
# Left accent strip
strip = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.25), Inches(7.5))
strip.fill.solid(); strip.fill.fore_color.rgb = BLUE; strip.line.fill.background()
# Gradient-ish accent rectangle bottom
bot = sl.shapes.add_shape(1, Inches(0), Inches(6.4), Inches(13.33), Inches(1.1))
bot.fill.solid(); bot.fill.fore_color.rgb = CORAL; bot.line.fill.background()

t1 = tb(sl, 0.6, 1.5, 12.0, 2.0)
rn(t1.text_frame, "CAUCHY INTEGRAL THEOREM", TF, 48, True, BLUE, PP_ALIGN.LEFT, True)

t2 = tb(sl, 0.6, 3.5, 12.0, 0.7)
rn(t2.text_frame, "Complex Analysis  |  Mathematics", BF, 22, False, BLACK, PP_ALIGN.LEFT, True)

t3 = tb(sl, 0.6, 4.3, 12.0, 1.5)
tf3 = t3.text_frame; tf3.word_wrap = True
for i,line in enumerate([
    "Presented by: [Your Name]",
    "SRM Institute of Science and Technology",
    "Course: Complex Analysis"
]):
    rn(tf3, line, BF, 18, False, BLACK, PP_ALIGN.LEFT, i==0)

# ══════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════
sl = new_slide()
heading_bar(sl, "AGENDA")

# Two columns
left = [
    "01  Introduction",
    "02  Why Complex Integration?",
    "03  Basic Definitions",
    "04  Cauchy Integral Theorem",
    "05  Conditions & Proof",
]
right = [
    "06  Cauchy Integral Formula",
    "07  Applications",
    "08  Worked Examples",
    "09  Conclusion & References",
]

tl = tb(sl, 0.4, 1.4, 6.2, 5.8)
tf = tl.text_frame; tf.word_wrap = True
first = True
for item in left:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False; p.space_before = Pt(10); p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = item; r.font.name = BF
    r.font.size = Pt(18); r.font.color.rgb = BLUE; r.font.bold = True

tr_ = tb(sl, 6.8, 1.4, 6.2, 5.8)
tf2 = tr_.text_frame; tf2.word_wrap = True
first = True
for item in right:
    p = tf2.paragraphs[0] if first else tf2.add_paragraph()
    first = False; p.space_before = Pt(10); p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = item; r.font.name = BF
    r.font.size = Pt(18); r.font.color.rgb = CORAL; r.font.bold = True

# ══════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION
# ══════════════════════════════════════════════
sl = new_slide()
heading_bar(sl, "Introduction to Complex Analysis")
bullet_box(sl, [
    "Complex analysis studies functions of complex variables:  z = x + iy",
    "Analytic (holomorphic) functions are differentiable w.r.t. complex variables — they have remarkable properties like infinite differentiability.",
    "Integration occurs along paths called contours in the complex plane, not just the real line.",
    "Augustin-Louis Cauchy discovered the central theorem of this field.",
    "Applications: Fluid dynamics · Electromagnetic theory · Quantum physics · Engineering mathematics",
], y=1.35)

# ══════════════════════════════════════════════
# SLIDE 4 — BASIC DEFINITIONS
# ══════════════════════════════════════════════
sl = new_slide()
heading_bar(sl, "Basic Definitions")

# Two-column layout: definition cards
defs = [
    ("Complex Function", "f(z) = u(x,y) + i·v(x,y),  where z = x+iy"),
    ("Analytic Function", "Differentiable at every point in a region. Examples: z², eᶻ, sin(z)"),
    ("Contour", "Smooth curve/path in the complex plane along which integration is performed"),
    ("Closed Contour", "Start point = End point (e.g. a circle). Critical for Cauchy's theorem"),
    ("Simply Connected", "Any closed contour can be shrunk to a point without leaving the region"),
    ("Derivative", "f′(z) = lim[Δz→0]  [f(z+Δz) − f(z)] / Δz"),
]
col1 = defs[:3]; col2 = defs[3:]

for col, ox in [(col1, 0.4), (col2, 6.9)]:
    for j, (name, desc) in enumerate(col):
        yy = 1.4 + j * 1.95
        box = sl.shapes.add_shape(1, Inches(ox), Inches(yy), Inches(6.1), Inches(1.75))
        box.fill.solid(); box.fill.fore_color.rgb = LTBLUE; box.line.fill.background()
        t = tb(sl, ox+0.12, yy+0.08, 5.9, 1.6)
        tf = t.text_frame; tf.word_wrap = True
        rn(tf, name, TF, 16, True, BLUE, PP_ALIGN.LEFT, True)
        rn(tf, desc, BF, 14, False, BLACK, PP_ALIGN.LEFT, False)

# ══════════════════════════════════════════════
# SLIDE 5 — THEOREM STATEMENT
# ══════════════════════════════════════════════
sl = new_slide()
heading_bar(sl, "Cauchy Integral Theorem — Statement")

t_intro = tb(sl, 0.4, 1.35, 12.5, 0.7)
rn(t_intro.text_frame,
   "If f(z) is analytic inside and on a simple closed contour C in a simply connected region, then:",
   BF, 17, False, BLACK, PP_ALIGN.LEFT, True)

formula_box(sl, ["∮C  f(z) dz  =  0"], x=2.5, y=2.1, w=8.0, h=1.2)

bullet_box(sl, [
    "The integral depends only on whether f(z) is analytic — NOT on the shape of the contour.",
    "If there are NO singularities inside C, the integral is always zero.",
    "This is analogous to conservative vector fields in real calculus.",
    "Developed by French mathematician Augustin-Louis Cauchy (1825).",
], y=3.45, size=17)

# ══════════════════════════════════════════════
# SLIDE 6 — CONDITIONS & PROOF
# ══════════════════════════════════════════════
sl = new_slide()
heading_bar(sl, "Conditions & Proof Outline")

# Conditions column
t_cl = tb(sl, 0.4, 1.35, 5.8, 5.9)
rn(t_cl.text_frame, "CONDITIONS", TF, 20, True, CORAL, PP_ALIGN.LEFT, True)
for cond in [
    "f(z) must be analytic inside the region",
    "f(z) analytic on the boundary (contour C)",
    "Contour C must be closed",
    "Region must be simply connected",
    "No singularities inside C",
]:
    p = t_cl.text_frame.add_paragraph()
    p.space_before = Pt(8); p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "✓  " + cond
    r.font.name = BF; r.font.size = Pt(16); r.font.color.rgb = BLACK

# Proof column
t_pr = tb(sl, 6.8, 1.35, 6.2, 5.9)
rn(t_pr.text_frame, "PROOF SKETCH", TF, 20, True, BLUE, PP_ALIGN.LEFT, True)
for step in [
    "Write f(z)=u+iv, dz=dx+idy",
    "Expand: ∮C f(z)dz = ∮C(udx−vdy) + i∮C(vdx+udy)",
    "Apply Green's Theorem → double integrals over R",
    "Use Cauchy–Riemann:  ∂u/∂x = ∂v/∂y  and  ∂u/∂y = −∂v/∂x",
    "Both double integrals = 0  ∴  ∮C f(z)dz = 0  ✓",
]:
    p = t_pr.text_frame.add_paragraph()
    p.space_before = Pt(7); p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "▸  " + step
    r.font.name = BF; r.font.size = Pt(15); r.font.color.rgb = BLACK

# ══════════════════════════════════════════════
# SLIDE 7 — CAUCHY INTEGRAL FORMULA
# ══════════════════════════════════════════════
sl = new_slide()
heading_bar(sl, "Cauchy Integral Formula")

rn(tb(sl,0.4,1.35,12.5,0.6).text_frame,
   "If f(z) is analytic inside and on C, and 'a' is any point inside C:",
   BF, 17, False, BLACK, PP_ALIGN.LEFT, True)

formula_box(sl, [
    "f(a)  =  (1 / 2πi) · ∮C  f(z) / (z − a) dz",
    "f⁽ⁿ⁾(a)  =  (n! / 2πi) · ∮C  f(z) / (z − a)ⁿ⁺¹ dz",
], x=1.0, y=2.0, w=11.2, h=1.6)

bullet_box(sl, [
    "The value of f at any interior point 'a' is fully determined by its values on the boundary C.",
    "Analytic functions are infinitely differentiable — every order of derivative exists.",
    "Leads to Taylor series expansions for analytic functions.",
    "Foundation for the Residue Theorem — the most powerful tool for evaluating integrals.",
], y=3.75, size=17)

# ══════════════════════════════════════════════
# SLIDE 8 — APPLICATIONS
# ══════════════════════════════════════════════
sl = new_slide()
heading_bar(sl, "Applications")

apps = [
    ("Path Independence", "∫C₁ f(z)dz = ∫C₂ f(z)dz  for any two paths between same endpoints."),
    ("Evaluating Integrals", "Analytic functions give ∮C f(z)dz = 0, simplifying many hard integrals."),
    ("Residue Theorem", "Extends the theorem to functions with singularities; used in signal processing."),
    ("Physics & Engineering", "Fluid dynamics, EM theory, quantum mechanics, control systems."),
]
for j, (name, desc) in enumerate(apps):
    yy = 1.38 + j * 1.5
    box = sl.shapes.add_shape(1, Inches(0.4), Inches(yy), Inches(12.5), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = LTBLUE if j % 2 == 0 else RGBColor(0xFF,0xF0,0xEE)
    box.line.fill.background()
    t = tb(sl, 0.55, yy+0.08, 12.0, 1.15)
    tf = t.text_frame; tf.word_wrap = True
    rn(tf, name, TF, 18, True, BLUE if j%2==0 else CORAL, PP_ALIGN.LEFT, True)
    rn(tf, desc, BF, 15, False, BLACK, PP_ALIGN.LEFT, False)

# ══════════════════════════════════════════════
# SLIDE 9 — WORKED EXAMPLES
# ══════════════════════════════════════════════
sl = new_slide()
heading_bar(sl, "Worked Examples")

examples = [
    ("Example 1", "∮C z² dz", "f(z)=z² is a polynomial — analytic everywhere in ℂ  →  Result = 0"),
    ("Example 2", "∮C (3z+2) dz", "f(z)=3z+2 is linear — analytic everywhere in ℂ  →  Result = 0"),
    ("Example 3", "∮C eᶻ dz", "f(z)=eᶻ is exponential — analytic everywhere in ℂ  →  Result = 0"),
]
for j, (label, formula, explanation) in enumerate(examples):
    yy = 1.38 + j * 1.9
    # Label
    t_lbl = tb(sl, 0.4, yy, 1.8, 0.5)
    rn(t_lbl.text_frame, label, TF, 17, True, CORAL, PP_ALIGN.LEFT, True)
    # Formula box
    fbox = tb(sl, 2.3, yy, 3.5, 0.5)
    fbox.fill.solid(); fbox.fill.fore_color.rgb = LTBLUE
    rn(fbox.text_frame, formula, MF, 17, True, BLUE, PP_ALIGN.CENTER, True)
    # Arrow
    arr = tb(sl, 5.9, yy, 0.6, 0.5)
    rn(arr.text_frame, "→", BF, 20, True, CORAL, PP_ALIGN.CENTER, True)
    # Explanation
    t_exp = tb(sl, 6.6, yy, 6.4, 1.3)
    tf = t_exp.text_frame; tf.word_wrap = True
    rn(tf, explanation, BF, 15, False, BLACK, PP_ALIGN.LEFT, True)

# ══════════════════════════════════════════════
# SLIDE 10 — CONCLUSION
# ══════════════════════════════════════════════
sl = new_slide()
heading_bar(sl, "Conclusion")

formula_box(sl, ["∮C f(z) dz = 0   (when f(z) is analytic inside and on C)"],
            x=1.0, y=1.35, w=11.2, h=0.9)

bullet_box(sl, [
    "Cauchy's theorem proves analytic functions integrate to zero over any closed contour.",
    "The theorem leads to the Cauchy Integral Formula, Residue Theorem, Taylor & Laurent series.",
    "Path independence: the integral depends only on analyticity, not the contour shape.",
    "Essential across mathematics, physics, and engineering — one of the most elegant results in all of mathematics.",
], y=2.4, size=17)

# ══════════════════════════════════════════════
# SLIDE 11 — REFERENCES
# ══════════════════════════════════════════════
sl = new_slide()
heading_bar(sl, "References")
bullet_box(sl, [
    "Kreyszig, E. — Advanced Engineering Mathematics (10th Ed.)",
    "Churchill, R.V. & Brown, J.W. — Complex Variables and Applications (9th Ed.)",
    "Class lecture notes on Complex Analysis",
    "Ahlfors, L.V. — Complex Analysis (3rd Ed.), McGraw-Hill",
], y=1.5, size=19, numbered=True)

# ══════════════════════════════════════════════
# SLIDE 12 — THANK YOU
# ══════════════════════════════════════════════
sl = new_slide(DARKBG)
strip2 = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.25), Inches(7.5))
strip2.fill.solid(); strip2.fill.fore_color.rgb = CORAL; strip2.line.fill.background()

t_ty = tb(sl, 1.5, 2.2, 10.5, 1.8)
rn(t_ty.text_frame, "Thank You", TF, 72, True, WHITE, PP_ALIGN.CENTER, True)

t_sub = tb(sl, 1.5, 4.2, 10.5, 0.7)
rn(t_sub.text_frame, "Questions & Discussion Welcome", BF, 22, False, CORAL, PP_ALIGN.CENTER, True)

t_nm = tb(sl, 1.5, 5.1, 10.5, 0.6)
rn(t_nm.text_frame, "Complex Analysis  —  SRM Institute of Science and Technology", BF, 16, False, WHITE, PP_ALIGN.CENTER, True)

# ── SAVE ──
OUT = r"c:\Users\ayyub\.gemini\antigravity\scratch\NewProject\Cauchy_Integral_Theorem.pptx"
prs.save(OUT)
print(f"Saved ({len(prs.slides)} slides): {OUT}")
